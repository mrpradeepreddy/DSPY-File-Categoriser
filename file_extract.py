import os
import fitz                      # PDFs
import docx                      # Word
import pandas as pd              # Excel/CSV/ODS
from pptx import Presentation    # PowerPoint
from odf import text, teletype
import pytesseract               # OCR for images
from PIL import Image
from odf.opendocument import load
import pypandoc


def parse_pdf_text(file_path: str) -> str:
    """
    Extract text/content from many common file types.
    (For images/audio returns a message instead of text.)
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    # --- Text-based files ---
    if ext == '.pdf':
        doc = fitz.open(file_path)
        full_text = "".join(page.get_text() for page in doc)
        doc.close()
        return full_text.strip()

    elif ext in ['.txt', '.py', '.js', '.json', '.xml', '.yaml', '.yml', '.html', '.css', '.md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    elif ext in ['.docx', '.doc']:
        document = docx.Document(file_path)
        return "\n".join(p.text for p in document.paragraphs)

    elif ext == '.odt':
        odt_doc = load(file_path)
        allparas = odt_doc.getElementsByType(text.P)
        return "\n".join(teletype.extractText(p) for p in allparas)

    elif ext == '.rtf':
        return pypandoc.convert_text(open(file_path, encoding='utf-8').read(),'plain', format='rtf', extra_args=['--standalone'])

    # --- Spreadsheet files ---
    elif ext in ['.csv', '.xls', '.xlsx', '.ods']:
        df = pd.read_excel(file_path) if ext != '.csv' else pd.read_csv(file_path)
        return df.to_string(index=False)

    # --- Presentation files ---
    elif ext in ['.pptx', '.ppt']:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_runs.append(paragraph.text)
        return "\n".join(text_runs)

    elif ext == '.odp':
        # ODP = OpenDocument Presentation
        od_doc = load(file_path)
        allparas = od_doc.getElementsByType(text.P)
        return "\n".join(teletype.extractText(p) for p in allparas)

    # --- Image files (need OCR) --- 
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.svg', '.tiff']:
        try:
            image = Image.open(file_path)        # open the image
            text = pytesseract.image_to_string(image)  # run OCR
            return text.strip()
        except Exception as e:
            print(f"Error reading image file {file_path}: {e}")
            return None

    # --- Audio files (need speech-to-text) --- not implemented here
    elif ext in ['.mp3', '.wav', '.aac', '.flac']:
        return "(Audio file detected – run speech-to-text here if you need text.)"

    else:
        return f"(Unsupported or binary file type {ext})"
